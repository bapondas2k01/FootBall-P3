from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_COLOR = RGBColor(31, 78, 121)  # Dark blue
ACCENT_COLOR = RGBColor(0, 176, 80)   # Green
TEXT_COLOR = RGBColor(51, 51, 51)     # Dark gray
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    # Add title text
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_before = Pt(10)
    p.space_after = Pt(10)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.line_spacing = 1.3
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "Football-P3", "Interactive Arcade Soccer Game\nBuilt with Phaser 3 & TypeScript")

# Slide 2: Project Overview
add_content_slide(prs, "Project Overview", [
    "🎮 Dynamic football (soccer) game with intelligent AI opponents",
    "📊 Developer: BaponDas | Student ID: 223071014",
    "🎯 Status: Development Complete & Fully Integrated",
    "⚙️ Game Engine: Phaser 3 (v3.87.0)",
    "💾 Tech Stack: TypeScript, Vite, JavaScript",
    "🎨 157+ Game Assets Fully Integrated",
    "📱 Cross-Platform Web-Based Game"
])

# Slide 3: Core Features
add_content_slide(prs, "Core Game Features", [
    "✨ Interactive Gameplay with smooth controls & realistic physics",
    "🤖 AI Opponent System with adaptive behavior",
    "🎭 Multiple Game Scenes (Landing, Mode Select, Loading, Game, Victory)",
    "📊 Real-time HUD (Score, Timer, Status Indicators)",
    "🎵 Complete Audio System (2 themes + 9 sound effects)",
    "⚡ Optimized Performance (60 FPS, Vite optimization)",
    "🎮 Responsive Controls (Arrow Keys/WASD + Space for kick)"
])

# Slide 4: Game Modes & Controls
add_content_slide(prs, "Game Modes & Controls", [
    "👥 Single Player Mode: Play against intelligent AI opponent",
    "👫 Two-Player Mode: Local multiplayer gameplay",
    "⌨️ Movement: Arrow Keys or WASD for multi-directional control",
    "🦵 Ball Interaction: Space bar to kick with physics-based trajectory",
    "🛡️ Defensive: Slide tackles for defensive gameplay",
    "⚡ Real-time Keyboard Input Handling",
    "🎯 Game Duration: 90 seconds per match"
])

# Slide 5: AI & Physics System
add_content_slide(prs, "AI & Physics System", [
    "🧠 Intelligent AI Controller: Advanced opponent behavior",
    "📍 Adaptive AI: Responds to player & ball position",
    "🎚️ Difficulty Levels: Varying AI behavior tuning",
    "⚽ Ball Physics: Realistic bouncing & collision detection",
    "👥 Player Collision: Player-to-player & boundary handling",
    "🌍 Gravity System: Optional physics simulation (arcade style)",
    "📊 Velocity Tracking: Visual debugging capabilities"
])

# Slide 6: Visual Assets
add_content_slide(prs, "Visual Assets (68 Images)", [
    "👤 Player Characters: Messi, Ronaldo, Blue Team, Red Team styles",
    "⚽ Soccer Ball: Multiple sprites with variants",
    "🥅 Goal Posts: 6 different styles (standard, cartoon, fixed, variants)",
    "🏟️ Backgrounds: Stadium view, clean field, tileset",
    "🎮 UI Elements: Buttons (play, restart, back), banners, panels",
    "✨ Special Effects: Stun stars, goal text overlay",
    "🎬 73 Animation Frames: 2 sequences (37 frames each @ 10fps)"
])

# Slide 7: Audio System
add_content_slide(prs, "Audio System (15 Files)", [
    "🎵 Background Music: 2 unique game themes",
    "🔊 Sound Effects:",
    "   • Ball kick, bounce, goal cheer, post hit",
    "   • UI button clicks, slide tackles",
    "   • Game start fanfare, victory fanfare",
    "   • Referee whistle",
    "🎚️ Integrated Audio Management",
    "🔔 Sound Feedback for all game interactions"
])

# Slide 8: Project Structure
add_content_slide(prs, "Project Architecture", [
    "📁 Source Organization:",
    "   • main.ts - Entry point",
    "   • Player.ts, Ball.ts, Goal.ts - Game mechanics",
    "   • AIController.ts - AI logic",
    "   • GameUI.ts - UI management",
    "🎬 Scene Flow: LoadingScene → LandingScene → ModeSelectScene →",
    "             StartScene → GameScene → VictoryScene",
    "⚙️ Configuration: gameConfig.json for settings"
])

# Slide 9: Technical Stack
add_content_slide(prs, "Technology Stack", [
    "🎮 Game Engine: Phaser 3.87.0",
    "📝 Language: TypeScript (Type-safe development)",
    "🔧 Build Tool: Vite (Fast dev & production builds)",
    "📦 Plugins: phaser3-rex-plugins (Extended functionality)",
    "📐 Resolution: 1152×768 pixels (optimized for web)",
    "⚡ Performance: 60 FPS stable performance",
    "🎨 Pixel Art Mode: Enabled for retro aesthetic"
])

# Slide 10: Asset Integration
add_content_slide(prs, "Asset Management", [
    "📦 Total Assets: 157 items organized & integrated",
    "🖼️ Images: 68 files in public/assets/images/",
    "🎵 Audio: 15 files in public/assets/audio/",
    "🎬 Animations: 73 frames in public/assets/animations/",
    "🔤 Fonts: Retro Pixel Arcade font",
    "✅ All assets fully integrated into game engine",
    "📋 Master configuration: asset-pack.json"
])

# Slide 11: Getting Started
add_content_slide(prs, "Quick Start Guide", [
    "1️⃣ Install Dependencies: npm install",
    "2️⃣ Run Development Server: npm run dev (http://localhost:8080)",
    "3️⃣ Build for Production: npm run build",
    "4️⃣ Preview Build: npm run preview",
    "📚 Available Commands:",
    "   • npm run download-assets - Download external assets",
    "   • npm run integrate-assets - Integrate local assets"
])

# Slide 12: Key Specifications
add_content_slide(prs, "Game Specifications", [
    "🎮 Screen Resolution: 1152 × 768 pixels",
    "⚡ Frame Rate: 60 FPS (stable)",
    "⏱️ Game Duration: 90 seconds per match",
    "🎉 Goal Celebration Time: 1.5 seconds",
    "👥 Score Display: Real-time HUD updates",
    "🏅 Victory Screen: Game outcome display with replay",
    "🔊 Total Audio Files: 15 (2 themes + 13 effects)"
])

# Slide 13: Development Highlights
add_content_slide(prs, "Development Highlights", [
    "✅ Full TypeScript Implementation: Type-safe development",
    "✅ Clean Architecture: Modular scene-based system",
    "✅ Physics Engine: Realistic collision & movement",
    "✅ AI System: Adaptive opponent behavior",
    "✅ Asset Optimization: 157 assets optimized & integrated",
    "✅ Performance Optimized: Chunk optimization with Vite",
    "✅ Production Ready: Build & deployment ready"
])

# Slide 14: Conclusion
add_title_slide(prs, "Ready for Deployment", "Football-P3: A Complete Arcade Soccer Gaming Experience")

# Save presentation
prs.save("Football_P3_Presentation.pptx")
print("✅ Presentation created: Football_P3_Presentation.pptx")
